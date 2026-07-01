from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.core.config import settings
from src.models.report import Report
from src.models.user import User

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


@dataclass
class Theme:
    key: str
    background: str
    accent: str
    font_heading: str
    font_body: str


THEMES: dict[str, Theme] = {
    "boardroom": Theme("boardroom", "#0a0f1e", "#2563eb", "Arial", "Arial"),
    "consulting": Theme("consulting", "#ffffff", "#1a1a2e", "Calibri", "Calibri"),
    "startup": Theme("startup", "#0a0a0a", "#f97316", "Calibri", "Calibri"),
    "editorial": Theme("editorial", "#fafaf8", "#7c3aed", "Georgia", "Georgia"),
    "academic": Theme("academic", "#ffffff", "#374151", "Times New Roman", "Times New Roman"),
}


def _rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def _text_color_for_bg(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    return RGBColor(0x11, 0x11, 0x11) if luminance > 0.6 else RGBColor(0xFF, 0xFF, 0xFF)


def _output_dir(report_id: str) -> Path:
    out_dir = Path(settings.upload_dir) / "reports" / report_id
    out_dir.mkdir(parents=True, exist_ok=True)
    return out_dir


class _SlideBuilder:
    def __init__(self, prs: Presentation, theme: Theme, user: User, is_branded: bool):
        self.prs = prs
        self.theme = theme
        self.user = user
        self.is_branded = is_branded
        self.accent = user.brand_primary if (is_branded and user.brand_primary) else theme.accent
        self.font_heading = user.brand_font if (is_branded and user.brand_font) else theme.font_heading
        self.text_color = _text_color_for_bg(theme.background)

    def _new_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank layout
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb(self.theme.background)
        if self.is_branded and self.user.brand_logo_url:
            self._add_logo_placeholder(slide)
        return slide

    def _add_logo_placeholder(self, slide):
        # Real logo image bytes aren't available locally (brand_logo_url is a
        # remote URL); reserve the corner with the brand name as a stand-in.
        box = slide.shapes.add_textbox(SLIDE_W - Inches(2.2), Inches(0.2), Inches(2), Inches(0.4))
        tf = box.text_frame
        tf.text = "BRANDED"
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(9)
        run.font.color.rgb = _rgb(self.accent)

    def _add_text(self, slide, text, left, top, width, height, size, bold=False, color=None, align=PP_ALIGN.LEFT, font=None):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font or self.theme.font_body
        run.font.color.rgb = color or self.text_color
        return box

    def title_slide(self, hook: str):
        slide = self._new_slide()
        self._add_text(
            slide, hook, Inches(1), Inches(2.6), Inches(11.3), Inches(2.5),
            size=40, bold=True, color=self.text_color, font=self.font_heading,
        )
        self._add_text(
            slide, "A DataBrief Story", Inches(1), Inches(5.4), Inches(6), Inches(0.5),
            size=14, color=_rgb(self.accent),
        )

    def about_data_slide(self, context: str):
        slide = self._new_slide()
        self._add_text(slide, "About This Data", Inches(1), Inches(0.7), Inches(10), Inches(0.8),
                        size=28, bold=True, color=_rgb(self.accent), font=self.font_heading)
        self._add_text(slide, context, Inches(1), Inches(2), Inches(11), Inches(3),
                        size=18, color=self.text_color)

    def finding_slide(self, index: int, finding: dict[str, Any]):
        slide = self._new_slide()
        self._add_text(slide, f"Finding {index}", Inches(1), Inches(0.6), Inches(6), Inches(0.6),
                        size=14, color=_rgb(self.accent))
        self._add_text(slide, finding.get("description", ""), Inches(1), Inches(1.4), Inches(7), Inches(3),
                        size=22, bold=True, color=self.text_color, font=self.font_heading)
        value = finding.get("value")
        if value is not None:
            self._add_text(slide, f"{value:,.2f}" if isinstance(value, (int, float)) else str(value),
                            Inches(8.5), Inches(2.2), Inches(3.5), Inches(1.5),
                            size=48, bold=True, color=_rgb(self.accent), align=PP_ALIGN.CENTER)

    def statement_slide(self, eyebrow: str, statement: str, big: bool = False):
        slide = self._new_slide()
        self._add_text(slide, eyebrow, Inches(1), Inches(0.7), Inches(10), Inches(0.6),
                        size=14, color=_rgb(self.accent))
        self._add_text(
            slide, statement, Inches(1), Inches(2.4), Inches(11.3), Inches(3),
            size=44 if big else 24, bold=big, color=self.text_color, font=self.font_heading,
            align=PP_ALIGN.CENTER if big else PP_ALIGN.LEFT,
        )

    def appendix_slide(self, findings: list[dict[str, Any]]):
        slide = self._new_slide()
        self._add_text(slide, "Appendix — Raw Stats", Inches(1), Inches(0.6), Inches(10), Inches(0.6),
                        size=24, bold=True, color=_rgb(self.accent), font=self.font_heading)
        y = 1.5
        for f in findings:
            self._add_text(
                slide, f"[{f.get('type')}] {f.get('description')}", Inches(1), Inches(y), Inches(11), Inches(0.5),
                size=12, color=self.text_color,
            )
            y += 0.5


def build_pptx(report: Report, user: User, story_arc: dict[str, Any]) -> str:
    theme_key = report.pptx_theme or "boardroom"
    theme = THEMES.get(theme_key, THEMES["boardroom"])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builder = _SlideBuilder(prs, theme, user, report.is_branded)

    builder.title_slide(story_arc.get("hook", "Your Data Story"))
    builder.about_data_slide(story_arc.get("context", ""))

    findings = story_arc.get("findings") or []
    for i, finding in enumerate(findings[:5], start=1):
        builder.finding_slide(i, finding)

    climax = story_arc.get("climax")
    if climax:
        builder.statement_slide("THE MAIN INSIGHT", climax.get("description", ""), big=True)

    if story_arc.get("implication"):
        builder.statement_slide("WHAT THIS MEANS", story_arc["implication"])
    if story_arc.get("action"):
        builder.statement_slide("WHAT TO DO NEXT", story_arc["action"])
    if story_arc.get("open_question"):
        builder.statement_slide("THE QUESTION TO INVESTIGATE", story_arc["open_question"])

    if findings:
        builder.appendix_slide(findings)

    from src.services import storage_service
    import tempfile, os
    if storage_service._r2_configured():
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        tmp.close()
        prs.save(tmp.name)
        object_key = f"reports/{report.id}/report.pptx"
        storage_service.upload_file(tmp.name, object_key, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        os.remove(tmp.name)
        return object_key
    out_path = _output_dir(str(report.id)) / "report.pptx"
    prs.save(str(out_path))
    return str(out_path)
